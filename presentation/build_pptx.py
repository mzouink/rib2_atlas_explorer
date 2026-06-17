#!/usr/bin/env python
"""Build the Filters & Metrics deck from analysis.json + thumbs + matrix.png. Run in base env."""
import json
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRES = os.path.dirname(os.path.abspath(__file__))
A = json.load(open(f"{PRES}/analysis.json"))
N = A["n"]
NAVY = RGBColor(0x1F, 0x32, 0x4D); GREY = RGBColor(0x5b, 0x66, 0x70)
GREEN = RGBColor(0x2E, 0x7D, 0x32); RED = RGBColor(0xC0, 0x40, 0x2C)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
B = prs.slide_layouts[6]


def tb(s, l, t, w, h):
    x = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); x.text_frame.word_wrap = True
    return x.text_frame


def run(p, txt, sz, color=NAVY, bold=False, italic=False):
    r = p.add_run(); r.text = txt; r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; return r


def img(s, path, l, t, mw, mh):
    """mw, mh in pixels; l, t in inches. Scales to fit, anchored at (l, t)."""
    if not os.path.exists(path):
        return
    w, h = Image.open(path).size
    r = min(mw / w, mh / h)
    s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w * r / 96), Inches(h * r / 96))


# ---- title ----
s = prs.slides.add_slide(B)
tf = tb(s, 0.8, 2.4, 11.7, 2.5)
run(tf.paragraphs[0], "RNA Atlas Explorer — Filters & Metrics", 36, NAVY, True)
run(tf.add_paragraph(), f"What each selection criterion means, an example it keeps vs one it drops, "
    f"and how much of the {N:,}-fold curated set each removes.", 18, GREY)
run(tf.add_paragraph(), "Live: ddc01lh56i5th.cloudfront.net", 14, RGBColor(0x2e,0x6f,0x95))


def card(s, x, title, color, e, val_label):
    """good/bad example: chemmap-colored 3D + props (top), seq·chemmap·motif tracks (bottom)."""
    head = tb(s, x, 1.45, 6.0, 0.4)
    run(head.paragraphs[0], title, 17, color, True)
    if not e:
        run(tb(s, x, 2.0, 6, 0.5).paragraphs[0], "(none)", 14, GREY); return
    img(s, f"{PRES}/thumbs/{e['id']}.png", x, 1.85, 2.5 * 96, 1.95 * 96)
    tf = tb(s, x + 2.55, 1.85, 3.45, 2.1); tf.word_wrap = True
    run(tf.paragraphs[0], f"{val_label}: {e['value']}", 15, color, True)
    rows = [e["id"], e["name"] or "—", f"{e['length']} nt · pLDDT {e['plddt']:.0f}",
            f"novelty TM1 {e['best_tm1']:.3f}" if e["best_tm1"] is not None else "novelty —",
            f"nearest {e['near']}" + (f" — {e['near_title'][:38]}" if e["near_title"] else ""),
            f"SHAPE agr {e['shape_agr']}" if e["shape_agr"] is not None else "SHAPE agr —",
            f"compact {e['contact_ratio']} · paired {e['bp_fraction']} · {e['n_tert']} tert"]
    for txt in rows:
        run(tf.add_paragraph(), txt, 10.5, NAVY)
    img(s, f"{PRES}/tracks/{e['id']}.png", x, 4.0, 6.0 * 96, 3.0 * 96)


for m in A["metrics"]:
    s = prs.slides.add_slide(B)
    tf = tb(s, 0.5, 0.3, 12.3, 1.1)
    run(tf.paragraphs[0], m["name"], 28, NAVY, True)
    run(tf.add_paragraph(), m["blurb"], 14, GREY)
    p = tf.add_paragraph()
    run(p, f"Drops {m['drop']:,} of {N:,} folds  ({m['drop_pct']}%)", 15, RED, True)
    run(p, f"   ·   keeps {m['keep']:,}", 13, GREEN)
    card(s, 0.5, "KEEP — strong example", GREEN, m["good"], "value")
    card(s, 6.9, "DROP — removed by this filter", RED, m["bad"], "value")

# ---- combination truth-table ----
s = prs.slides.add_slide(B)
run(tb(s, 0.5, 0.3, 12.3, 0.7).paragraphs[0], "Filter combinations — what survives", 26, NAVY, True)
img(s, f"{PRES}/combos.png", 0.4, 0.95, 7.4 * 96, 6.4 * 96)
tf = tb(s, 8.3, 1.3, 4.7, 5.8)
run(tf.paragraphs[0], "How to read", 16, NAVY, True)
for t in ["Each row = one on/off combination of the 6 core filters (✓ = applied).",
          "All 64 combinations; top row = no filters (all 7,757).",
          "“% left” / “# left” = folds passing every applied filter (AND).",
          "Sorted by # kept — tightest combinations at the bottom.",
          "Filters overlap, so stacking them removes less than the sum of each alone.",
          "Full table (incl. all 8 filters): combinations.csv."]:
    run(tf.add_paragraph(), "• " + t, 12.5, GREY)

# ---- scaling / time estimate ----
s = prs.slides.add_slide(B)
run(tb(s, 0.5, 0.3, 12.3, 0.7).paragraphs[0], "Scaling to the full predicted atlas — cost & bottleneck", 24, NAVY, True)
run(tb(s, 0.5, 0.95, 12.3, 0.6).paragraphs[0],
    "Everything here is computed for the 7,757 curated folds. Running the same pipeline atlas-wide is gated by one step.", 13, GREY)
img(s, f"{PRES}/scaling.png", 0.6, 1.55, 8.4 * 96, 4.6 * 96)
tf = tb(s, 9.1, 1.7, 3.9, 5.4)
run(tf.paragraphs[0], "Bottleneck", 16, RED, True)
run(tf.add_paragraph(), "USalign novelty: ~0.5–3 min/fold. Every other step is seconds or sub-second.", 12.5, NAVY)
run(tf.add_paragraph(), "Feasible", 15, GREEN, True)
run(tf.add_paragraph(), "414k high-confidence set: ~1–2 days on a ~200-slot LSF array.", 12.5, NAVY)
run(tf.add_paragraph(), "Impractical", 15, RED, True)
run(tf.add_paragraph(), "Full ~23M atlas: novelty alone is ~2–3 months on a 200-slot array — and most are low-confidence. Even the cheap steps become week-scale.", 12.5, NAVY)
run(tf.add_paragraph(), "Recommended", 15, NAVY, True)
run(tf.add_paragraph(), "Run the full pipeline over the 414k strict index; skip novelty on the low-confidence rest (or throw far more slots at it).", 12.5, NAVY)

H = json.load(open(f"{PRES}/hvt.json"))


def exline(tf, e):
    p = tf.add_paragraph()
    run(p, f"{e['id']}  ", 11, NAVY, True)
    run(p, f"{e['length']} nt · {e['name'][:38]}  ·  tm {e['best_tm1']:.2f} · agr {e['shape_agr']} · "
        f"comp {e['contact_ratio']} · {e['n_tert']}t{' · PK' if e['pseudoknot'] else ''}", 10.5, GREY)


# ---- slide 1: the binding constraint is SIZE ----
s = prs.slides.add_slide(B)
run(tb(s, 0.5, 0.35, 12.3, 0.7).paragraphs[0], "Target selection collapses on size", 28, NAVY, True)
img(s, f"{PRES}/funnel.png", 0.5, 1.2, 7.3 * 96, 3.7 * 96)
tf = tb(s, 8.1, 1.2, 4.9, 5.6)
run(tf.paragraphs[0], "The hard reality", 16, NAVY, True)
for t in ["Un-assisted single-particle cryo-EM of RNA wants ~150+ nt (~50 kDa).",
          "Ribonanza-2 design regions cap ~130 nt — so every candidate is small for cryo-EM.",
          f"Of {H['funnel'][0][1]:,} novel folds, only {H['funnel'][2][1]} are ≥100 nt and just "
          f"{H['funnel'][3][1]} clears the full quality gate."]:
    run(tf.add_paragraph(), "• " + t, 13, GREY)
run(tf.add_paragraph(), "→ The decision that sets the list:", 14, RED, True)
for t in ["Accept Fab/scaffold-assisted cryo-EM? If yes → hundreds of small but striking folds open up.",
          "If no → we’re limited to the ~12 largest, fully de-risked folds."]:
    run(tf.add_paragraph(), "• " + t, 13, NAVY)

# ---- slide 2: two concrete routes (with structure + sequence/chemmap visuals) ----
s = prs.slides.add_slide(B)
run(tb(s, 0.5, 0.25, 12.3, 0.6).paragraphs[0], "Proposed list — two routes (rank, then curate together)", 23, NAVY, True)


def example_viz(s, x, e):
    """3D (chemmap+motif) + seq·chemmap tracks + a label, in a ~3.0in column."""
    img(s, f"{PRES}/thumbs/{e['id']}.png", x, 1.95, 1.7 * 96, 1.35 * 96)
    img(s, f"{PRES}/tracks/{e['id']}.png", x + 1.55, 1.95, 1.55 * 96, 1.35 * 96)
    tf = tb(s, x, 3.2, 3.05, 1.0); tf.word_wrap = True
    run(tf.paragraphs[0], f"{e['id']} · {e['length']} nt", 10.5, NAVY, True)
    run(tf.add_paragraph(), (e["name"] or "")[:42], 9.5, GREY)
    run(tf.add_paragraph(), f"tm {e['best_tm1']:.2f} · agr {e['shape_agr']} · comp {e['contact_ratio']} · "
        f"{e['n_tert']}t{' · PK' if e['pseudoknot'] else ''}", 9.5, GREY)


def route(s, x, title, color, crit, sub, examples):
    tf = tb(s, x, 0.95, 6.2, 0.95); tf.word_wrap = True
    run(tf.paragraphs[0], title, 15, color, True)
    run(tf.add_paragraph(), crit, 11, NAVY)
    run(tf.add_paragraph(), sub, 10, GREY, italic=True)
    for i, e in enumerate(examples[:2]):
        example_viz(s, x + i * 3.15, e)


route(s, 0.45, f"Route A · Largest, fully de-risked  ({H['route1_n']})", GREEN,
      "len ≥ 80 · pLDDT ≥ 88 · SHAPE-supported · compact ≥ 0.5 · novel (TM1 ≤ 0.40) · ≥1 tertiary",
      "Best shot at a usable map with least scaffolding; mostly RNAMake assemblies ~90–100 nt.", H["route1"])
route(s, 6.85, f"Route B · Most striking architecture  ({H['route2_n']})", RGBColor(0x6a, 0x4c, 0x93),
      "novel · SHAPE-supported · pseudoknot · rare tertiary motif",
      "Smaller → cryo-EM needs a chaperone/scaffold, but high scientific interest.", H["route2"])

# remaining examples as compact lines
tf = tb(s, 0.45, 4.55, 6.2, 2.0); tf.word_wrap = True
run(tf.paragraphs[0], "More Route A candidates:", 11, GREEN, True)
for e in H["route1"][2:]:
    exline(tf, e)
tf = tb(s, 6.85, 4.55, 6.2, 2.0); tf.word_wrap = True
run(tf.paragraphs[0], "More Route B candidates:", 11, RGBColor(0x6a, 0x4c, 0x93), True)
for e in H["route2"][2:]:
    exline(tf, e)
run(tb(s, 0.45, 6.9, 12.4, 0.5).paragraphs[0],
    "3D colored by 2A3 reactivity (blue protected→red reactive) + motif sticks; tracks = sequence · DMS · 2A3 · motifs.  "
    "Mostly synthetic designs — large+SHAPE+novel natural pool is near-empty. Pull ~5–8/route → ~10–15 portfolio.",
    10, GREY, italic=True)

# ---- slide 3: what de-risks a slot + the missing axis ----
s = prs.slides.add_slide(B)
run(tb(s, 0.5, 0.35, 12.3, 0.7).paragraphs[0], "What de-risks a cryo-EM slot — and the axis we’re missing", 23, NAVY, True)
tf = tb(s, 0.6, 1.2, 12.1, 3.6); tf.word_wrap = True
run(tf.paragraphs[0], "cryo-EM is expensive — each metric buys a piece of confidence before we commit:", 14, GREY)
for head, body in [
    ("Model says it folds", "pLDDT + gPDE."),
    ("Experiment says it folds (orthogonal!)", "SHAPE agreement — guards against a confidently-wrong model."),
    ("It’s a compact single particle", "compactness (C1′ contact ratio) + base-paired fraction."),
    ("It’s worth solving", "novelty vs the PDB + the nearest-entry title as a sanity check."),
    ("It’s interesting", "tertiary motifs, pseudoknots."),
]:
    p = tf.add_paragraph(); run(p, "• " + head + " — ", 13.5, NAVY, True); run(p, body, 13, RGBColor(0x34, 0x40, 0x4d))
tf = tb(s, 0.6, 5.0, 12.1, 2.1); tf.word_wrap = True
run(tf.paragraphs[0], "The axis we don’t yet measure: rigidity / single dominant conformation", 16, RED, True)
run(tf.add_paragraph(), "A floppy, multi-conformer RNA gives no usable map regardless of how it scores above. We have one "
    "prediction per fold, so we can’t see conformational spread.", 13, NAVY)
run(tf.add_paragraph(), "Proposal: before committing, re-predict each shortlisted fold with N seeds and keep the low inter-seed-RMSD "
    "(single rigid fold) ones — cheap for a ~10–30 shortlist.", 13, NAVY)

# ---- discussion questions ----
s = prs.slides.add_slide(B)
run(tb(s, 0.5, 0.35, 12.3, 0.8).paragraphs[0], "Open questions for the discussion", 28, NAVY, True)
tf = tb(s, 0.7, 1.5, 11.9, 5.2)
for i, q in enumerate([
    "What size can the cryo-EM team realistically attempt — and are they open to Fab / scaffold assistance for < ~150 nt?",
    "RNA-only targets, or RNA–protein complexes?",
    "Natural (F–H, more reliably folded) vs synthetic (A–E, more novel) — and what mix?",
    "Hard-require experimental SHAPE support, or allow model-only candidates?",
    "How many targets are we scoping — ~5 hero targets, or a ~20–30 shortlist?",
], 1):
    p = tf.add_paragraph(); run(p, f"{i}.  ", 16, RGBColor(0x2e, 0x6f, 0x95), True); run(p, q, 15, NAVY)
run(tf.add_paragraph(), " ", 8)
run(tf.add_paragraph(), "Proposal: bring the explorer to the meeting, tune thresholds live, and walk the top of the ranked list together.", 13, GREY, italic=True)

out = f"{PRES}/filters_and_metrics.pptx"
prs.save(out)
print("wrote", out, "slides:", len(prs.slides._sldIdLst))
